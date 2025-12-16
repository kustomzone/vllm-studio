"""Artifact generators for creating various file types."""
from pathlib import Path
from typing import Any, Dict, List, Optional, Union
import uuid
import json


def create_file(
    content: Union[str, bytes],
    filename: str,
    artifacts_dir: Path,
    mime_type: str = "text/plain",
) -> Dict[str, Any]:
    """Create a generic file artifact."""
    artifact_id = str(uuid.uuid4())[:8]
    safe_filename = f"{artifact_id}_{filename}"
    filepath = artifacts_dir / safe_filename
    
    if isinstance(content, str):
        filepath.write_text(content)
    else:
        filepath.write_bytes(content)
    
    return {
        "id": artifact_id,
        "type": "file",
        "filename": filename,
        "path": str(filepath),
        "mime_type": mime_type,
        "size": filepath.stat().st_size,
    }


def create_spreadsheet(
    data: List[List[Any]],
    filename: str = "spreadsheet.xlsx",
    artifacts_dir: Path = None,
    sheet_name: str = "Sheet1",
    headers: Optional[List[str]] = None,
) -> Dict[str, Any]:
    """Create an Excel spreadsheet artifact."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
    except ImportError:
        return {"error": "openpyxl not installed"}
    
    if artifacts_dir is None:
        artifacts_dir = Path("/tmp/vllm_artifacts")
        artifacts_dir.mkdir(exist_ok=True)
    
    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name
    
    row_offset = 1
    
    # Add headers if provided
    if headers:
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
            cell.font = Font(bold=True, color="FFFFFF")
        row_offset = 2
    
    # Add data
    for row_idx, row_data in enumerate(data, row_offset):
        for col_idx, value in enumerate(row_data, 1):
            ws.cell(row=row_idx, column=col_idx, value=value)
    
    # Auto-adjust column widths
    for column_cells in ws.columns:
        length = max(len(str(cell.value or "")) for cell in column_cells)
        ws.column_dimensions[column_cells[0].column_letter].width = min(length + 2, 50)
    
    artifact_id = str(uuid.uuid4())[:8]
    safe_filename = f"{artifact_id}_{filename}"
    filepath = artifacts_dir / safe_filename
    wb.save(filepath)
    
    return {
        "id": artifact_id,
        "type": "spreadsheet",
        "filename": filename,
        "path": str(filepath),
        "mime_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "size": filepath.stat().st_size,
        "rows": len(data) + (1 if headers else 0),
        "columns": len(data[0]) if data else 0,
    }


def create_pdf(
    content: str,
    filename: str = "document.pdf",
    artifacts_dir: Path = None,
    title: Optional[str] = None,
) -> Dict[str, Any]:
    """Create a PDF document artifact."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.units import inch
    except ImportError:
        return {"error": "reportlab not installed"}
    
    if artifacts_dir is None:
        artifacts_dir = Path("/tmp/vllm_artifacts")
        artifacts_dir.mkdir(exist_ok=True)
    
    artifact_id = str(uuid.uuid4())[:8]
    safe_filename = f"{artifact_id}_{filename}"
    filepath = artifacts_dir / safe_filename
    
    doc = SimpleDocTemplate(str(filepath), pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    # Add title if provided
    if title:
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=30,
        )
        story.append(Paragraph(title, title_style))
    
    # Add content paragraphs
    for para in content.split('\n\n'):
        if para.strip():
            story.append(Paragraph(para.strip(), styles['Normal']))
            story.append(Spacer(1, 12))
    
    doc.build(story)
    
    return {
        "id": artifact_id,
        "type": "pdf",
        "filename": filename,
        "path": str(filepath),
        "mime_type": "application/pdf",
        "size": filepath.stat().st_size,
    }


def create_presentation(
    slides: List[Dict[str, Any]],
    filename: str = "presentation.pptx",
    artifacts_dir: Path = None,
) -> Dict[str, Any]:
    """Create a PowerPoint presentation artifact.
    
    slides format:
    [
        {"title": "Slide Title", "content": "Bullet point text", "layout": "title"},
        {"title": "Content Slide", "bullets": ["Point 1", "Point 2"], "layout": "bullets"},
    ]
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return {"error": "python-pptx not installed"}
    
    if artifacts_dir is None:
        artifacts_dir = Path("/tmp/vllm_artifacts")
        artifacts_dir.mkdir(exist_ok=True)
    
    prs = Presentation()
    
    for slide_data in slides:
        layout_type = slide_data.get("layout", "bullets")
        
        if layout_type == "title":
            layout = prs.slide_layouts[0]  # Title slide
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_data.get("title", "")
            if slide_data.get("subtitle"):
                slide.placeholders[1].text = slide_data["subtitle"]
        else:
            layout = prs.slide_layouts[1]  # Title and content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_data.get("title", "")
            
            if slide_data.get("bullets"):
                body = slide.placeholders[1]
                tf = body.text_frame
                for i, bullet in enumerate(slide_data["bullets"]):
                    if i == 0:
                        tf.text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
            elif slide_data.get("content"):
                body = slide.placeholders[1]
                body.text = slide_data["content"]
    
    artifact_id = str(uuid.uuid4())[:8]
    safe_filename = f"{artifact_id}_{filename}"
    filepath = artifacts_dir / safe_filename
    prs.save(filepath)
    
    return {
        "id": artifact_id,
        "type": "presentation",
        "filename": filename,
        "path": str(filepath),
        "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "size": filepath.stat().st_size,
        "slides": len(slides),
    }
